## So now we take the TEP model from first_tep.py where we really jsut used the defautl model times 3 
# and we add predictions. we ant to see what happens when a tap comes in at certain delays after the pulse
#-------------------------------------
#Init
#-------------------------------------
        

OG_HNN=False


## Original HNN
from neuron import h
h.nrn_load_dll('C:\\Users\\ckohl\\Miniconda3\\Lib\\site-packages\\hnn_core\\nrnmech.dll')
import os.path as op
if OG_HNN==True:
    #import hnn_core
    from hnn_core import simulate_dipole, Params, Network, read_params
else:
    #import hnn_core_ca
    from hnn_core_ca import simulate_dipole, Params, Network, read_params
import json
import numpy as np
import matplotlib.pyplot as plt
import sys        
sys.path.append('C:\\Users\\ckohl\\Documents\\HNN core code\\')
import my_hnn_core_functions as m
from pptx import Presentation #https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
from pptx.util import Inches

saving_dpl=False
my_param_out_dir='C:\\Users\\ckohl\\hnn_out\\'
dump_dir="C:\\Users\\ckohl\\Desktop\\Current\\HNN Core\\"     


#-------------------------------------
#PPT init
#-------------------------------------
left = 0
top=Inches(0.25)
height = Inches(7)
width=Inches(10)
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "NEW (r) fits"
subtitle.text = "plot_from_hnn.py"

# time delay
delays=[0,20,30,40,50,60,70,80,100,120,140,160,180,200,220,240,250,300,400,500,600]#[50, 100, 150, 200, 300, 400, 500, 600]
delayed_input_type=['default','Supra_ERPYesSupraT']# also try triple later?

## pick a param file and tripple those inputs
param_files=['default','Supra_ERPYesSupraT']
add10=True
triple_pulse=True
    
    
keep_dpls=[]
for p in param_files:
    for delayed_p in delayed_input_type:
        for delay_t in delays:
            dict_param=m.get_dict_param(my_param_out_dir,p)
            base_params=[]
            these_params=[]
            base_params = Params(dict_param)
            these_params=Params(dict_param)

            #adjust params for sarah's model (don't have a base for supra)
            if OG_HNN==False:
                #find where sarah and default are different
                param_to_comp=['thesis_erp_100%','default']
                dict_param=m.get_dict_param(my_param_out_dir,param_to_comp[0])
                comp1=Params(dict_param)
                dict_param=m.get_dict_param(my_param_out_dir,param_to_comp[1])
                comp2=Params(dict_param)
                  
                paramoi=[]
                for pt in comp1:
                    if comp1[pt]!= comp2[pt]:
                        paramoi.append(pt)
                
                for pt in paramoi:
                    these_params[pt]=comp1[pt]

            if triple_pulse==True:
                # triple them
                Inputs=[]
                prox_count=0
                dist_count=0
                input_times=[]
                for pi in these_params:
                    if pi[0:8]=='t_evprox':
                        prox_count=prox_count+1
                        Inputs.append('Prox'+str(prox_count))
                        input_times.append(these_params[pi])
                    if pi[0:8]=='t_evdist':
                        dist_count=dist_count+1
                        Inputs.append('Dist'+str(dist_count))
                        input_times.append(these_params[pi])
                #sort inputs
                sort_i=np.argsort(input_times)         
                zipped_pairs = zip(sort_i,Inputs)       
                Inputs = [x for _, x in sorted(zipped_pairs)]        
                     
                mylabel=['Spikes','Seed','t','sd','2Pyr_AMPA','2Pyr_NMDA','2Bask_AMPA','2Bas_NMDA','5Pyr_AMPA','5Pyr_NMDA','5Bask_AMPA','5Bask_NMDA']
                actuallabel_pre=['numspikes','prng_seedcore','t','sigma_t','gbar','gbar','gbar','gbar','gbar','gbar','gbar','gbar']
                actuallabel_post=['','','','','_L2Pyr_ampa','_L2Pyr_nmda','_L2Basket_ampa','_L2Basket_nmda','_L5Pyr_ampa','_L5Pyr_nmda','_L5Basket_ampa','_L5Basket_nmda']
                time_delay=20
                for input in range(0,len(Inputs)):
                    #lets go through inputs and copy them
                    if Inputs[input][0]=='P':
                        it_through_these=mylabel
                        actuallabel_mid='_evprox_'
                        input_number=prox_count
                    else:
                        it_through_these=mylabel[0:10]
                        actuallabel_mid='_evdist_'
                        input_number=dist_count
                        
                    for input_p in range(0,len(it_through_these)):    
                            
                        this_param=actuallabel_pre[input_p]+actuallabel_mid+Inputs[input][-1]+actuallabel_post[input_p]
                        #this param with new number
                        this_new_param1=actuallabel_pre[input_p]+actuallabel_mid+str(int(Inputs[input][-1])+input_number)+actuallabel_post[input_p]
                        this_new_param2=actuallabel_pre[input_p]+actuallabel_mid+str(int(Inputs[input][-1])+input_number+1)+actuallabel_post[input_p]
                        #make 2 copies of ths param and add them to the param dict (with new input numbers) (unless its time, then change it)
                            
                        if mylabel[input_p]=='t':
                            these_params[this_new_param1]=these_params[this_param]+time_delay
                            these_params[this_new_param2]=these_params[this_param]+time_delay*2
                        else:
                            these_params[this_new_param1]=these_params[this_param]
                            these_params[this_new_param2]=these_params[this_param]
                        
                    if Inputs[input][0]=='P':
                        prox_count=prox_count+1
                    else:
                        dist_count=dist_count+1
                    
            keep_these_params=these_params        
                    
                    
            #add delayed input
            dict_param=m.get_dict_param(my_param_out_dir,delayed_p)
            add_params=[]
            add_params = Params(dict_param)
            
            #adjust params for sarah's model (don't have a base for supra)
            if OG_HNN==False:
                #find where sarah and default are different
                param_to_comp=['thesis_erp_100%','default']
                dict_param=m.get_dict_param(my_param_out_dir,param_to_comp[0])
                comp1=Params(dict_param)
                dict_param=m.get_dict_param(my_param_out_dir,param_to_comp[1])
                comp2=Params(dict_param)
                  
                paramoi=[]
                for pt in comp1:
                    if comp1[pt]!= comp2[pt]:
                        paramoi.append(pt)
                
                for pt in paramoi:
                    add_params[pt]=comp1[pt]

            Inputs=[]
            prox_count=0
            dist_count=0
            input_times=[]
            for pi in these_params:
                if pi[0:8]=='t_evprox':
                    prox_count=prox_count+1
                    Inputs.append('Prox'+str(prox_count))
                    input_times.append(these_params[pi])
                if pi[0:8]=='t_evdist':
                    dist_count=dist_count+1
                    Inputs.append('Dist'+str(dist_count))
                    input_times.append(these_params[pi])
                    
            New_Inputs=[]
            new_prox_count=prox_count
            new_dist_count=dist_count
            input_times=[]
            for pi in add_params:
                if pi[0:8]=='t_evprox':
                    new_prox_count=new_prox_count+1
                    New_Inputs.append('Prox'+str(new_prox_count))
                    input_times.append(these_params[pi])
                if pi[0:8]=='t_evdist':
                    new_dist_count=new_dist_count+1
                    New_Inputs.append('Dist'+str(new_dist_count))
                    input_times.append(these_params[pi])
            #sort inputs
            sort_i=np.argsort(input_times)         
            zipped_pairs = zip(sort_i,New_Inputs)       
            New_Inputs = [x for _, x in sorted(zipped_pairs)]        
                 
            mylabel=['Spikes','Seed','t','sd','2Pyr_AMPA','2Pyr_NMDA','2Bask_AMPA','2Bas_NMDA','5Pyr_AMPA','5Pyr_NMDA','5Bask_AMPA','5Bask_NMDA']
            actuallabel_pre=['numspikes','prng_seedcore','t','sigma_t','gbar','gbar','gbar','gbar','gbar','gbar','gbar','gbar']
            actuallabel_post=['','','','','_L2Pyr_ampa','_L2Pyr_nmda','_L2Basket_ampa','_L2Basket_nmda','_L5Pyr_ampa','_L5Pyr_nmda','_L5Basket_ampa','_L5Basket_nmda']
            for input in range(0,len(New_Inputs)):
                #lets go through inputs and copy them
                if New_Inputs[input][0]=='P':
                    it_through_these=mylabel
                    actuallabel_mid='_evprox_'
                    input_number=prox_count
                else:
                    it_through_these=mylabel[0:10]
                    actuallabel_mid='_evdist_'
                    input_number=dist_count
                    
                for input_p in range(0,len(it_through_these)):    
              
                        
                    this_param=actuallabel_pre[input_p]+actuallabel_mid+str(int(New_Inputs[input][-1])-input_number)+actuallabel_post[input_p]
                    #this param with new number
                    this_new_param1=actuallabel_pre[input_p]+actuallabel_mid+New_Inputs[input][-1]+actuallabel_post[input_p]
                    #make 2 copies of ths param and add them to the param dict (with new input numbers) (unless its time, then change it)
                        
                    if mylabel[input_p]=='t':
                        these_params[this_new_param1]=add_params[this_param]+delay_t
                    else:
                        these_params[this_new_param1]=add_params[this_param]
                    
                # if Inputs[input][0]=='P':
                    # prox_count=prox_count+1
                # else:
                    # dist_count=dist_count+1
                    
                   
            
            #adjust time and scale                
            these_params['tstop']=delay_t+200
            these_params['dipole_scalefctr']=400
            keep_these_params['dipole_scalefctr']=400
            if add10==True:
                for pars in these_params:
                    if pars[0:2]=='t_':
                        these_params[pars]=these_params[pars]+10      

          
            #turn off smoothing
            these_params['dipole_smooth_win']=0
            # simulate
            net=[]
            dpls=[]
            net = Network(these_params)
            dpls = simulate_dipole(net, n_jobs=1, n_trials=1)
            m.core_output_basic(dpls,net,name,data,False,True, False)
            keep_dpls.append(dpls[0].dpl['agg'])
            #m.plot_hnn_core_output(dpls,net,'Supra_OG_x3',False,True,False,these_params)
            data='C:\\Users\\ckohl\\Desktop\\Current\\TEP\\Digitised from papers\\Avg0204_filt.txt'
            if OG_HNN==True:
                name='OG'
            else:
                name='SARAH'
            if p[0]=='d':
                name=name+'_'+p+'_x3'
            else:
                name=name+'_'+'supra'+'_x3'
            if add10==True:
                name=name+'_+10'
            name=name+'_+'+delayed_p+str(delay_t)
            m.core_output_basic(dpls,net,name,data,False,True, False)
            #m.plot_hnn_core_output(dpls,net,p,True,False,True,param_oi_t,these_params,base_params)
            plt.pause(.1)
            plt.savefig(dump_dir+'\\temp.png')
            img_path = dump_dir+'\\temp.png'
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(img_path, left, top, height=height,width=width)
            plt.close()

prs.save(dump_dir+'\\new_pred_smooth_test.pptx')



#plot dpls like rob
# fig = plt.figure()
# for delay in range(0,len(delays)):
    # if delays[delay]<50:
        # colour='green'
    # elif delays[delay]<400:
        # colour='red'
    # else:
        # colour='blue'
    # plt.plot(keep_dpls[delay],color=colour)
    
    
    
# fig = plt.figure()
# for delay in range(0,len(delays)):
    # if delays[delay]<50:
        # colour='green'
    # elif delays[delay]<400:
        # colour='red'
    # else:
        # colour='blue'
    # plt.plot(keep_dpls[delay],color=colour)
    # datalen=len(keep_dpls[delay])

    # dump_dir="C:\\Users\\ckohl\\Desktop\\Current\\HNN Core\\"
    # import pickle
    # f = open(dump_dir+p+'_dpl.txt', 'wb')
    # pickle.dump(dpls, f)
    # f.close()
    # f = open(dump_dir+p+'_net.txt', 'wb')
    # pickle.dump(net, f)
    # f.close()